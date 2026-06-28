# rag/loader.py

from bs4 import BeautifulSoup
from pathlib import Path

from langchain_core.documents import Document
from langchain_community.document_loaders import (
    PyPDFLoader,
    TextLoader,
)


# ==========================================================
# HTML loader
# ==========================================================
def load_html(file_path):
    with open(
        file_path,
        "r",
        encoding="utf-8",
        errors="ignore",
    ) as f:
        html = f.read()

    soup = BeautifulSoup(html, "lxml")

    for tag in soup(["script", "style", "nav", "footer"]):
        tag.decompose()

    title = (
        soup.title.text.strip()
        if soup.title
        else str(file_path)
    )

    body = soup.body
    if body:
        text = body.get_text(separator="\n")
    else:
        text = soup.get_text(separator="\n")

    text = "\n".join(
        line.strip()
        for line in text.split("\n")
        if line.strip()
    )

    return Document(
        page_content=text,
        metadata={
            "source": str(file_path),
            "title": title,
            "type": "html",
        },
    )


# ==========================================================
# PDF loader
# ==========================================================
def load_pdf(file_path):
    loader = PyPDFLoader(str(file_path))
    docs = loader.load()

    for doc in docs:
        doc.metadata["source"] = str(file_path)
        doc.metadata["type"] = "pdf"

    return docs


# ==========================================================
# TXT loader
# ==========================================================
def load_txt(file_path):
    loader = TextLoader(
        str(file_path),
        encoding="utf-8",
    )
    docs = loader.load()

    for doc in docs:
        doc.metadata["source"] = str(file_path)
        doc.metadata["type"] = "txt"

    return docs


# ==========================================================
# DOCX loader
# ==========================================================
def load_docx(file_path):
    import docx

    doc = docx.Document(str(file_path))

    text = "\n".join(
        para.text
        for para in doc.paragraphs
        if para.text.strip()
    )

    title = Path(file_path).stem

    return [Document(
        page_content=text,
        metadata={
            "source": str(file_path),
            "title": title,
            "type": "docx",
        },
    )]


# ==========================================================
# PPTX loader
# ==========================================================
def load_pptx(file_path):
    from pptx import Presentation

    prs = Presentation(str(file_path))
    docs = []

    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())

        if not texts:
            continue

        text = "\n".join(texts)
        title = Path(file_path).stem + f" - slide {i}"

        docs.append(Document(
            page_content=text,
            metadata={
                "source": str(file_path),
                "title": title,
                "type": "pptx",
                "slide": i,
            },
        ))

    return docs


# ==========================================================
# Generic loader
# ==========================================================
def load_file(file_path):
    suffix = file_path.suffix.lower()

    if suffix == ".pdf":
        return load_pdf(file_path)

    if suffix in (".html", ".htm"):
        return [load_html(file_path)]

    if suffix == ".txt":
        return load_txt(file_path)

    if suffix == ".docx":
        return load_docx(file_path)

    if suffix == ".pptx":
        return load_pptx(file_path)

    return []
